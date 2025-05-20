import pandas as pd
import requests
import fitz  # PyMuPDF
from bs4 import BeautifulSoup
import pytesseract
import platform
from pdf2image import convert_from_path, pdfinfo_from_path
from pathlib import Path
from PIL import Image
from pptx import Presentation
import textract
from langdetect import detect
import os, glob, shutil
import platform
from transformers import pipeline, AutoModelForCTC, AutoTokenizer
from pydub import AudioSegment
import tempfile

# English vocabulary for detecting poorly encoded PDFs
english_dict = pd.read_csv(
    "https://github.com/dwyl/english-words/raw/master/words_alpha.txt", header=None
)
english_dict = list(english_dict.iloc[:, 0].values)


def setup_directories(data_path):
    "setup requied and standardized directory structure"

    # creating base path
    if not os.path.isdir(data_path):
        os.makedirs(data_path)

    # path for raw documents, PDFs, HTML, etc.
    if not os.path.isdir(f"{data_path}raw_files/"):
        os.makedirs(f"{data_path}raw_files/")

    # path for .txt documents
    if not os.path.isdir(f"{data_path}txt_files/"):
        os.makedirs(f"{data_path}txt_files/")

    # path for potential transformed .txt files, e.g., stemmed, stopwords removed, etc.
    if not os.path.isdir(f"{data_path}transformed_txt_files/"):
        os.makedirs(f"{data_path}transformed_txt_files/")

    # path for transformed data CSVs
    if not os.path.isdir(f"{data_path}csv_outputs/"):
        os.makedirs(f"{data_path}csv_outputs/")

    # path for reports, plots, and outputs
    if not os.path.isdir(f"{data_path}visual_outputs/"):
        os.makedirs(f"{data_path}visual_outputs/")


def clear_directories(
    processor,
    raw_files=True,
    txt_files=True,
    transformed_txt_files=True,
    csv_outputs=True,
    visual_outputs=True,
):
    "clear the data in different directories"

    def clear_directory(dir_path):
        files = os.listdir(dir_path)
        for file in files:
            os.remove(f"{dir_path}{file}")

    if raw_files:
        clear_directory(f"{processor.data_path}raw_files/")
    if txt_files:
        clear_directory(f"{processor.data_path}txt_files/")
    if transformed_txt_files:
        clear_directory(f"{processor.data_path}transformed_txt_files/")
    if csv_outputs:
        clear_directory(f"{processor.data_path}csv_outputs/")
    if visual_outputs:
        clear_directory(f"{processor.data_path}visual_outputs/")


def generate_metadata_file(data_path, metadata_addt_column_names):
    "generate an initial metadata file. metadata_addt_column_names is list of additional columns necessary for the analysis (publication date, etc.)"

    # only create the file if it doesn't exist already
    if not (os.path.isfile(f"{data_path}metadata.csv")):
        metadata = pd.DataFrame(
            columns=[
                "text_id",
                "web_filepath",
                "local_raw_filepath",
                "local_txt_filepath",
                "detected_language",
            ]
            + metadata_addt_column_names
        )
        metadata.to_csv(f"{data_path}metadata.csv", index=False)
    else:
        metadata = pd.read_csv(f"{data_path}metadata.csv")
        metadata.text_id = range(
            1, len(metadata) + 1
        )  # ensure creation of a unique text id field
    return metadata


def download_document(metadata, data_path, text_id, web_filepath):
    "download a file from a URL and update the metadata file"

    if str(web_filepath) == "" or str(web_filepath) == "nan":
        return None
    else:
        web_filepath = web_filepath.split(",")[
            0
        ]  # may have multiple URLs stored in field, take only first (english)
        web_filepath = web_filepath.splitlines()[
            0
        ]  # may have multiple URLS stored on separate lines, take only first

        # first check if this file already downloaded
        if (
            not (os.path.isfile(f"{data_path}raw_files/{text_id}.txt"))
            and not (os.path.isfile(f"{data_path}raw_files/{text_id}.csv"))
            and not (os.path.isfile(f"{data_path}raw_files/{text_id}.doc"))
            and not (os.path.isfile(f"{data_path}raw_files/{text_id}.docx"))
            and not (os.path.isfile(f"{data_path}raw_files/{text_id}.html"))
            and not (os.path.isfile(f"{data_path}raw_files/{text_id}.pdf"))
            and not (os.path.isfile(f"{data_path}raw_files/{text_id}.txt"))
            and not (os.path.isfile(f"{data_path}raw_files/{text_id}.jpg"))
            and not (os.path.isfile(f"{data_path}raw_files/{text_id}.mp3"))
            and not (os.path.isfile(f"{data_path}raw_files/{text_id}.mp4"))
            and not (os.path.isfile(f"{data_path}raw_files/{text_id}.m4a"))
            and not (os.path.isfile(f"{data_path}raw_files/{text_id}.wav"))
        ):
            print(f"{data_path}raw_files/{text_id}")
            try:  # try downloading the file first
                headers = {
                    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.36"
                }
                response = requests.get(web_filepath, headers=headers)
                if response.status_code == 200:
                    print("Success")
                else:
                    print(f"Failed to access URL. Status code: {response.status_code}")

                content_type = response.headers.get("content-type")

                if "application/pdf" in content_type:
                    ext = ".pdf"
                elif "text/html" in content_type:
                    ext = ".html"
                elif "officedocument" in content_type:
                    ext = ".docx"
                elif "application/msword" in content_type:
                    ext = ".doc"
                elif "text/csv" in content_type:
                    ext = ".csv"
                elif "text/plain" in content_type:
                    ext = ".txt"
                elif "image/jpeg" in content_type:
                    ext = ".jpg"
                elif (
                    web_filepath[-3:] == "mp3"
                ):  # the mp3 content_type returns 'application/octet-stream'
                    ext = ".mp3"
                elif web_filepath[-3:] == "mp4":
                    ext = ".mp4"
                elif web_filepath[-3:] == "wav":
                    ext = ".wav"
                elif web_filepath[-3:] == "m4a":
                    ext = ".m4a"
                else:
                    ext = ""

                if ext != "":
                    file = open(f"{data_path}raw_files/{text_id}{ext}", "wb+")
                    file.write(response.content)
                    file.close()

                    metadata.loc[metadata.text_id == text_id, "local_raw_filepath"] = (
                        f"{data_path}raw_files/{text_id}{ext}"
                    )
                    return metadata
                else:
                    return None
            except:
                try:  # try just copying the local txt file if not a URL
                    if ".txt" in web_filepath:
                        shutil.copyfile(
                            web_filepath, f"{data_path}raw_files/{text_id}.txt"
                        )
                        metadata.loc[
                            metadata.text_id == text_id, "local_raw_filepath"
                        ] = f"{data_path}raw_files/{text_id}.txt"
                    else:
                        return None
                except:
                    return None


def parse_pdf(pdf_path):
    "parse a pdf and return text string"
    # Open the PDF file
    doc = fitz.open(pdf_path)
    return_text = ""

    # Iterate over each page in the PDF
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)  # Load each page
        return_text += "[newpage] " + page.get_text(
            "text"
        )  # Extract the text content of the page

    return return_text


def parse_ocr_pdf(
    data_path, pdf_path, windows_tesseract_path=None, windows_poppler_path=None
):
    """convert a scanned PDF to text and return text string. From https://www.geeksforgeeks.org/python-reading-contents-of-pdf-using-ocr-optical-character-recognition/
    parameters:
        :data_path: str: filepath where all files will be created and stored
        :pdf_path: str: path of the pdf file
        :windows_tesseract_path: str: path of tesseract .exe file (https://linuxhint.com/install-tesseract-windows/)
        :windows_poppler_path: str: path of poppler .exe file (https://blog.alivate.com.au/poppler-windows/)
    """
    if platform.system() == "Windows":
        pytesseract.pytesseract.tesseract_cmd = windows_tesseract_path
    # storing image files temporarily
    PDF_file = Path(pdf_path)
    image_file_list = []

    if platform.system() == "Windows":
        info = pdfinfo_from_path(
            pdf_path, userpw=None, poppler_path=windows_poppler_path
        )
        max_pages = info["Pages"]
        pdf_pages = []
        for i in range(1, max_pages + 1):
            print(f"OCR conversion (stage 1/3) of {pdf_path}, page {i} / {max_pages}")
            tmp_pdf_pages = convert_from_path(
                PDF_file,
                500,
                poppler_path=windows_poppler_path,
                first_page=i,
                last_page=i,
            )
            pdf_pages.append(tmp_pdf_pages[0])
    else:
        info = pdfinfo_from_path(pdf_path, userpw=None)
        max_pages = info["Pages"]
        pdf_pages = []
        for i in range(1, max_pages + 1):
            print(f"OCR conversion (stage 1/3) of {pdf_path}, page {i} / {max_pages}")
            tmp_pdf_pages = convert_from_path(PDF_file, 500, first_page=i, last_page=i)
            pdf_pages.append(tmp_pdf_pages[0])

    # iterate through pages
    counter = 1
    for page_enumeration, page in enumerate(pdf_pages, start=1):
        print(f"OCR conversion (stage 2/3) of {pdf_path}, page {counter} / {max_pages}")
        counter += 1

        # Create a file name to store the image
        filename = f"{data_path}/page_{page_enumeration:03}.jpg"

        # Save the image of the page in system
        page.save(filename, "JPEG")
        image_file_list.append(filename)

    # Iterate from 1 to total number of pages
    return_text = ""
    counter = 1
    for image_file in image_file_list:
        print(
            f"OCR conversion (stage 3/3) of {pdf_path}, page {counter} / {len(image_file_list)}"
        )
        counter += 1

        text = str(((pytesseract.image_to_string(Image.open(image_file)))))
        text = text.replace("-\n", "")
        return_text += "[newpage] " + text

    return return_text


def parse_html(html_path):
    "parse an html file and return text string"
    file = open(html_path, "r", encoding="UTF-8")
    html = file.read()
    soup = BeautifulSoup(html, features="html.parser")
    # kill all script and style elements
    for script in soup(["script", "style"]):
        script.extract()  # rip it out

    # get text
    return_text = soup.get_text()

    # break into lines and remove leading and trailing space on each
    lines = (line.strip() for line in return_text.splitlines())
    # break multi-headlines into a line each
    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
    # drop blank lines
    return_text = "\n".join(chunk for chunk in chunks if chunk)

    return return_text


def parse_word(doc_path):
    "parse a .docx or .doc file"
    return_text = (
        str(textract.process(doc_path))[2:-1].replace("\\n", "\n").replace("\\", "")
    )

    return return_text


def parse_pptx(doc_path):
    "parse a .pptx file"
    prs = Presentation(doc_path)
    text = ""

    for slide in prs.slides:
        text += "\n[newpage]\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


def parse_csv(doc_path):
    "parse a .csv into a markdown file"
    df = pd.read_csv(doc_path)
    return_text = df.to_markdown(index=False)

    return return_text


def parse_xlsx(doc_path):
    "parse a .csv into a markdown file"
    df = pd.read_excel(doc_path)
    return_text = df.to_markdown(index=False)

    return return_text


def parse_jpg(jpg_path):
    print("OCR conversion jpg")
    text = str(((pytesseract.image_to_string(Image.open(jpg_path)))))
    return_text = text.replace("-\n", "")
    return return_text


def load_asr(model_name):
    "load an automatic speech recognition model from huggingface"
    # Define the directory where the model will be saved
    save_directory = "./automatic-speech-recognition-models/" + model_name

    # Create the directory if it doesn't exist
    if not os.path.exists(save_directory):
        os.makedirs(save_directory)

    # Check if the model exists in the specified directory
    if os.listdir(save_directory):
        print(f"Loading model from {save_directory}")
        asr = pipeline("automatic-speech-recognition", model=save_directory)
    else:
        print(f"Model not found in {save_directory}, downloading and saving it.")
        asr = pipeline("automatic-speech-recognition", model=model_name)
        # and save it for future use
        asr.save_pretrained(save_directory)

    return asr


def parse_wav(wav_path, model_name="openai/whisper-base"):
    """
    Transcribes an .wav file into text using a specified model from Hugging Face and returns text.

    Parameters:
    - wav_path (str): Path to the .wav file to be transcribed.
    - model_name (str): Hugging Face model name for the STT task. Default is "openai/whisper-base".

    Returns:
    - Transcripted text from the wav
    """

    # if asr nor loaded then load it
    if "asr" not in globals():
        asr = load_asr(model_name=model_name)

    # Transcribe the WAV audio file directly
    return_text = asr(wav_path, compression_ratio_threshold=1.35)["text"]

    return return_text


def parse_mp3(raw_path, model_name="openai/whisper-base"):
    """
    Transcribes an .mp3 file into text using a specified model from Hugging Face and returns text.

    Parameters:
    - mp3_path (str): Path to the .mp3 file to be transcribed.
    - model_name (str): Hugging Face model name for the STT task. Default is "openai/whisper-base".

    Returns:
    - Transcripted text from the mp3
    """
    # Convert MP3 to WAV format using pydub
    if ".mp3" in raw_path:
        audio = AudioSegment.from_mp3(raw_path)
    elif ".mp4" in raw_path:
        audio = AudioSegment.from_file(raw_path, format="mp4")
    elif ".m4a" in raw_path:
        audio = AudioSegment.from_file(raw_path, format="m4a")

    # if asr nor loaded then load it
    if "asr" not in globals():
        asr = load_asr(model_name=model_name)

    # Use a temporary file to save the converted audio as WAV
    with tempfile.NamedTemporaryFile(suffix=".wav") as tmp_wav_file:
        audio.export(tmp_wav_file.name, format="wav")

        # Transcribe the WAV audio file directly
        return_text = asr(tmp_wav_file.name, return_timestamps=True)["text"]

    return return_text


def detect_language(stringx):
    "determine the language of a string"
    return detect(stringx)


def convert_to_text(
    metadata,
    data_path,
    text_id,
    windows_tesseract_path=None,
    windows_poppler_path=None,
    force_ocr=False,
):
    "convert a PDF or HTML file into raw text. In PDFs, new pages encoded with '[newpage]'"
    raw_path = metadata.loc[
        lambda x: x.text_id == text_id, "local_raw_filepath"
    ].values[0]

    # path exists (not nan) and is longer than 0
    raw_exists = type(raw_path) == str
    if raw_exists:
        raw_exists = raw_exists & (len(raw_path) > 0)

    if raw_exists:
        # first check if this file already converted
        if not (os.path.isfile(f"{data_path}txt_files/{text_id}.txt")):
            print(f"{data_path}txt_files/{text_id}.txt")
            # pdf file
            if ".pdf" in raw_path:
                try:
                    return_text = parse_pdf(raw_path)
                except:
                    return_text = ""

                # high proportion of non-english words could indicate poor encoding
                try:
                    detected_lang = detect(return_text)
                except:
                    detected_lang = "en"

                if detected_lang == "en":
                    alphas = set(
                        [x for x in return_text.lower().split(" ") if x.isalpha()]
                    )
                    eng_words = [x for x in alphas if x in english_dict]

                    try:
                        if len(alphas) == 0:
                            eng_dict = 0  # if eng_dict == 0, then ocr will be forced
                        else:
                            eng_dict = len(eng_words) / len(alphas)
                    except:
                        eng_dict = 1.0
                else:
                    eng_dict = 1.0

                try:
                    # check to see if ocr should be used
                    if (
                        len(return_text) == 0
                    ):  # this stops below condition throwing an error
                        use_ocr = True
                    else:
                        use_ocr = (
                            (
                                len(set(return_text.split("[newpage] "))) == 1
                            )  # if only empties, scan, needs to be OCR converted.
                            | (
                                (return_text.lower().count("/g") / len(return_text))
                                > 0.01
                            )  # If bunch of "/G"s, greater than 1% of all the characters, encoding error, like review of maritime transport 2006
                            | (return_text.lower().count("_") / len(return_text) > 0.05)
                            | (
                                return_text.lower().count("sqj") > 10
                            )  # if poorly digitized and a lot of 'sqj's
                            | (
                                return_text.lower().count("\x03") / len(return_text)
                                > 0.01
                            )  # if poorly digitized and a lot of '\x03's
                            | (
                                return_text.lower().count("\x01") / len(return_text)
                                > 0.01
                            )  # if poorly digitized and a lot of '\x01's
                            | (
                                return_text.lower().count("^") / len(return_text)
                                > 0.0001
                            )
                            | (
                                sum(
                                    [
                                        (
                                            1
                                            if return_text[i]
                                            == return_text[i - 1]
                                            == return_text[i - 2]
                                            and return_text[i].isalpha()
                                            else 0
                                        )
                                        for i in range(2, len(return_text))
                                    ]
                                )
                                / len(return_text)
                                > 0.0009
                            )  # many repeated letters is an error
                            | (
                                eng_dict < 0.8
                            )  # high proportion of out of vocabulary words
                            | (force_ocr)  # manually force OCR
                        )
                    if use_ocr:  # force OCR
                        return_text = parse_ocr_pdf(
                            data_path,
                            raw_path,
                            windows_tesseract_path,
                            windows_poppler_path,
                        )
                        # remove temporary image files from OCR
                        for f in glob.glob(f"{data_path}*.jpg"):
                            os.remove(f)
                except:
                    return_text = ""
            elif ".html" in raw_path:
                return_text = parse_html(raw_path)
            elif ".docx" in raw_path or ".doc" in raw_path:
                return_text = parse_word(raw_path)
            elif ".csv" in raw_path:
                return_text = parse_csv(raw_path)
            elif (".xlsx" in raw_path) or (".xls" in raw_path):
                return_text = parse_xlsx(raw_path)
            elif ".pptx" in raw_path:
                return_text = parse_pptx(raw_path)
            elif ".txt" in raw_path:
                file = open(f"{raw_path}", "r", encoding="UTF-8")
                return_text = file.read()
                file.close()
            elif ".jpg" in raw_path:
                return_text = parse_jpg(raw_path)
            elif ".mp3" in raw_path or ".mp4" in raw_path or ".m4a" in raw_path:
                try:
                    return_text = parse_mp3(raw_path=raw_path)
                except:
                    return_text = ""
            elif ".wav" in raw_path:
                try:
                    return_text = parse_wav(wav_path=raw_path)
                except:
                    return_text = ""

            # write text file
            file = open(f"{data_path}txt_files/{text_id}.txt", "wb+")
            file.write(return_text.encode("utf-8", "replace"))
            file.close()

            # update metadata file
            metadata.loc[lambda x: x.text_id == text_id, "local_txt_filepath"] = (
                f"{data_path}txt_files/{text_id}.txt"
            )
            try:
                detected_lang = detect(return_text)
            except:
                detected_lang = "en"
            metadata.loc[lambda x: x.text_id == text_id, "detected_language"] = (
                detected_lang
            )
            final_return = metadata
        else:
            final_return = None
    else:
        final_return = None

    return final_return


def refresh_local_metadata(metadata, data_path):
    """update the metadata to reflect the situation in the data path in terms of files/conversions done
    parameters:
        :metadata: pd.DataFrame: metadata dataframe of the object
        :data_path: str: filepath where all files are stored
    output:
        will update metadata for correct situation in terms of:
            -raw_files/
            -txt_files/
        and update downstream files (so that can't have a transformed txt or sentiment without the raw pdf, etc.),  in:
            -transformed_txt_files/
            -csv_outputs/
    """
    # raw path
    print(f"Syncing metadata to local file situation: checking raw paths...")

    def check_raw_exists(text_id):
        files = glob.glob(f"{data_path}raw_files/{text_id}.*")
        return files[0] if len(files) == 1 else ""

    raw_paths = [check_raw_exists(text_id) for text_id in metadata.text_id.values]
    metadata["local_raw_filepath"] = raw_paths

    # txt path
    print(f"Syncing metadata to local file situation: checking txt paths...")

    def check_txt_exists(text_id):
        txt_path = f"{data_path}txt_files/{text_id}.txt"
        return txt_path if os.path.exists(txt_path) else ""

    txt_paths = [check_txt_exists(text_id) for text_id in metadata.text_id.values]
    metadata["local_txt_filepath"] = txt_paths

    # remove transformed_txt/ files if there's not raw/txt path (old/stale file)
    print(
        f"Syncing metadata to local file situation: checking transformed_txt paths..."
    )

    def remove_transformed_files(text_id, txt_path):
        files = glob.glob(f"{data_path}transformed_txt_files/*_{text_id}.txt")
        if len(files) > 0:
            if (str(txt_path) == "") | (str(txt_path) == "nan"):
                for file in files:
                    os.remove(file)

    [
        remove_transformed_files(text_id, txt_path)
        for text_id, txt_path in dict(zip(metadata.text_id.values, txt_paths)).items()
    ]

    # csv_outputs/ files. Remove entries in CSVs where no raw/txt path exists
    print(f"Syncing metadata to local file situation: checking CSV outputs...")

    def remove_csv_outputs(text_id, txt_path, data):
        if (str(txt_path) == "") | (
            str(txt_path) == "nan"
        ):  # if no txt file, blank out all values for that in the csv
            data.loc[
                data.text_id == text_id, data.columns[data.columns != "text_id"]
            ] = ""
        return data

    csv_files = glob.glob(f"{data_path}csv_outputs/*.csv")
    counter = 1
    for file in csv_files:
        print(
            f"Syncing metadata to local file situation: checking CSV outputs: {counter}/{len(csv_files)}"
        )
        counter += 1

        data = pd.read_csv(file)
        for text_id, txt_path in dict(zip(metadata.text_id.values, txt_paths)).items():
            data = remove_csv_outputs(text_id, txt_path, data)
        data.to_csv(file, index=False)

    print(f"Syncing metadata to local file situation: Success!")

    return metadata
